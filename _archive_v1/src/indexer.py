import os
import json
import pickle
import time
import logging
from typing import List, Dict, Optional
from concurrent.futures import ThreadPoolExecutor, as_completed
import numpy as np
from sentence_transformers import SentenceTransformer
from langchain_text_splitters import RecursiveCharacterTextSplitter
from docx import Document
from pptx import Presentation
import openpyxl
from PyPDF2 import PdfReader
try:
    import win32com.client
    HAS_OUTLOOK = True
except ImportError:
    HAS_OUTLOOK = False
import pythoncom
from config import IndexingSettings

# Setup logger for indexer
logger = logging.getLogger(__name__)

class SimpleVectorStore:
    """Simple numpy-based vector store to replace ChromaDB."""
    
    def __init__(self, db_path: str = "./vector_db"):
        self.db_path = db_path
        self.documents: List[str] = []
        self.embeddings: Optional[np.ndarray] = None
        self.metadatas: List[Dict] = []
        self.ids: List[str] = []
        self.file_mtimes: Dict[str, float] = {}  # Track file modification times {filepath: mtime}
        
        # Ensure directory exists
        os.makedirs(db_path, exist_ok=True)
        
        # Load existing data if available
        self._load()
    
    def _load(self):
        """Load data from disk."""
        data_file = os.path.join(self.db_path, "vectors.pkl")
        if os.path.exists(data_file):
            try:
                with open(data_file, "rb") as f:
                    data = pickle.load(f)
                    self.documents = data.get("documents", [])
                    self.embeddings = data.get("embeddings")
                    self.metadatas = data.get("metadatas", [])
                    self.ids = data.get("ids", [])
                    self.file_mtimes = data.get("file_mtimes", {})
            except Exception as e:
                print(f"Warning: Could not load existing data: {e}")
    
    def _save(self):
        """Save data to disk."""
        data_file = os.path.join(self.db_path, "vectors.pkl")
        data = {
            "documents": self.documents,
            "embeddings": self.embeddings,
            "metadatas": self.metadatas,
            "ids": self.ids,
            "file_mtimes": self.file_mtimes
        }
        with open(data_file, "wb") as f:
            pickle.dump(data, f)
    
    def is_file_indexed(self, file_path: str) -> bool:
        """Check if file is already indexed with same modification time."""
        if file_path not in self.file_mtimes:
            return False
        try:
            current_mtime = os.path.getmtime(file_path)
            return self.file_mtimes[file_path] == current_mtime
        except OSError:
            return False
    
    def update_file_mtime(self, file_path: str):
        """Update the stored modification time for a file."""
        try:
            self.file_mtimes[file_path] = os.path.getmtime(file_path)
        except OSError:
            pass
    
    def upsert(self, documents: List[str], embeddings: List[List[float]], 
               metadatas: List[Dict], ids: List[str]):
        """Add or update documents."""
        # Remove existing documents with same IDs
        for new_id in ids:
            if new_id in self.ids:
                idx = self.ids.index(new_id)
                self.documents.pop(idx)
                self.metadatas.pop(idx)
                self.ids.pop(idx)
                if self.embeddings is not None:
                    self.embeddings = np.delete(self.embeddings, idx, axis=0)
        
        # Add new documents
        self.documents.extend(documents)
        self.metadatas.extend(metadatas)
        self.ids.extend(ids)
        
        # Add embeddings
        new_embeddings = np.array(embeddings)
        if self.embeddings is None or len(self.embeddings) == 0:
            self.embeddings = new_embeddings
        else:
            self.embeddings = np.vstack([self.embeddings, new_embeddings])
        
        self._save()
    
    def query(self, query_embeddings: List[List[float]], n_results: int = 5) -> Dict:
        """Query for similar documents using cosine similarity."""
        if self.embeddings is None or len(self.documents) == 0:
            return {
                'documents': [[]],
                'metadatas': [[]],
                'distances': [[]]
            }
        
        query_vec = np.array(query_embeddings[0])
        
        # Compute cosine similarity
        # Normalize vectors
        query_norm = query_vec / (np.linalg.norm(query_vec) + 1e-10)
        embeddings_norm = self.embeddings / (np.linalg.norm(self.embeddings, axis=1, keepdims=True) + 1e-10)
        
        # Cosine similarity
        similarities = np.dot(embeddings_norm, query_norm)
        
        # Get top n results
        n_results = min(n_results, len(self.documents))
        top_indices = np.argsort(similarities)[::-1][:n_results]
        
        # Convert similarity to distance (lower = better for consistency)
        distances = [1 - similarities[i] for i in top_indices]
        docs = [self.documents[i] for i in top_indices]
        metas = [self.metadatas[i] for i in top_indices]
        
        return {
            'documents': [docs],
            'metadatas': [metas],
            'distances': [distances]
        }
    
    def count(self) -> int:
        """Return number of documents."""
        return len(self.documents)
    
    def cleanup_deleted_files(self) -> int:
        """Remove entries for files that no longer exist on disk.
        Returns the number of entries removed."""
        if not self.metadatas:
            return 0
        
        # Find indices of entries to remove
        indices_to_remove = []
        for i, meta in enumerate(self.metadatas):
            source = meta.get('source', '')
            if isinstance(source, str) and source:
                # Check if it's a local file (not Outlook or other sources)
                if not source.startswith('outlook:') and not os.path.exists(source):
                    indices_to_remove.append(i)
        
        if not indices_to_remove:
            return 0
        
        # Remove entries in reverse order to maintain correct indices
        for idx in sorted(indices_to_remove, reverse=True):
            self.documents.pop(idx)
            self.metadatas.pop(idx)
            doc_id = self.ids.pop(idx)
            if self.embeddings is not None:
                self.embeddings = np.delete(self.embeddings, idx, axis=0)
            
            # Also remove from file_mtimes if present
            # Find the source path for this entry
            # (note: we already removed from metadatas, so we need to track removed sources)
        
        # Clean up file_mtimes for non-existent files
        stale_paths = [path for path in self.file_mtimes if not os.path.exists(path)]
        for path in stale_paths:
            del self.file_mtimes[path]
        
        self._save()
        return len(indices_to_remove)


class FileIndexer:
    def __init__(self, db_path: str = "./vector_db", model_name: str = "all-MiniLM-L6-v2"):
        self.errors = []  # Store errors for UI reporting
        self.model_name = model_name
        self.db_path = db_path

        # Use simple vector store instead of ChromaDB
        collection_name = model_name.replace("-", "_")
        full_db_path = os.path.join(db_path, collection_name)
        self.collection = SimpleVectorStore(db_path=full_db_path)
        
        # Use SentenceTransformer directly for embeddings
        self.embedding_model = SentenceTransformer(model_name)
        
        # Text splitter for chunking
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=100
        )
        
        print(f"Indexer initialized with {self.collection.count()} existing documents")

    def _get_embeddings(self, texts: List[str]) -> List[List[float]]:
        """Generate embeddings using SentenceTransformer."""
        embeddings = self.embedding_model.encode(texts, convert_to_numpy=True)
        return embeddings.tolist()

    def extract_text(self, file_path: str, settings: IndexingSettings = None) -> str:
        """Extract text from Office files based on extension with settings support."""
        ext = os.path.splitext(file_path)[1].lower()
        text = ""
        try:
            if ext == ".docx":
                doc = Document(file_path)
                text = "\n".join([p.text for p in doc.paragraphs])
            elif ext == ".pptx":
                prs = Presentation(file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
            elif ext == ".xlsx":
                text = self._extract_text_from_excel(file_path, settings)
            elif ext == ".txt" or ext == ".md":
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    text = f.read()
            elif ext == ".pdf":
                reader = PdfReader(file_path)
                for page in reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
        except Exception as e:
            error_msg = f"Error reading {file_path}: {e}"
            print(error_msg)
            self.errors.append({"file": file_path, "error": str(e)})
        return text
    
    def _extract_text_from_excel(self, file_path: str, settings: IndexingSettings = None) -> str:
        """Extract text from Excel file with settings-based limits."""
        wb = openpyxl.load_workbook(file_path, data_only=True)
        
        # Determine max rows and sheet filtering from settings
        if settings and settings.excel.limit_rows:
            max_rows = settings.excel.max_rows_per_sheet
        else:
            max_rows = None
        
        skip_raw_sheets = settings.excel.skip_raw_like_sheets if settings else True
        
        lines = []
        for sheet_name in wb.sheetnames:
            # Skip Raw/Log/History sheets if enabled
            if skip_raw_sheets:
                name_lower = sheet_name.lower()
                if any(x in name_lower for x in ['raw', 'log', 'history']):
                    print(f"  Skipping sheet '{sheet_name}' (Raw/Log/History)")
                    continue
            
            ws = wb[sheet_name]
            row_count = 0
            
            for row in ws.iter_rows(values_only=True):
                if max_rows is not None and row_count >= max_rows:
                    print(f"  Sheet '{sheet_name}' limited to {max_rows} rows")
                    break
                    
                values = [str(cell) for cell in row if cell is not None]
                if values:
                    lines.append(" ".join(values))
                row_count += 1
        
        return "\n".join(lines)

    def _should_skip_dir(self, dir_name: str, dir_path: str) -> bool:
        """Check if directory should be skipped."""
        # Skip hidden directories (starting with .)
        if dir_name.startswith('.'):
            return True
        
        # Skip common system/temp directories
        skip_dirs = {
            '$recycle.bin', 'appdata', 'programdata', 'windows',
            'program files', 'program files (x86)', '__pycache__',
            'node_modules', '.git', '.svn', 'venv', 'env', '.env',
            'temp', 'tmp', 'cache', '.cache', 'thumbs.db'
        }
        
        if dir_name.lower() in skip_dirs:
            return True
        
        return False

    def _should_skip_file(self, filename: str) -> bool:
        """Check if file should be skipped."""
        # Skip hidden files (starting with .)
        if filename.startswith('.'):
            return True
        
        # Skip temp Office files (starting with ~$)
        if filename.startswith('~$'):
            return True
        
        # Skip temp/backup files
        if filename.endswith('.tmp') or filename.endswith('.bak'):
            return True
        
        return False

    # 내용까지 인덱싱할 파일 확장자 (문서 파일)
    CONTENT_INDEXABLE_EXTS = {".docx", ".doc", ".pptx", ".ppt", ".xlsx", ".xls", ".txt", ".md", ".pdf"}
    
    # 완전히 제외할 파일 확장자
    SKIP_EXTENSIONS = {
        '.exe', '.dll', '.sys', '.drv', '.ocx',  # 실행/시스템
        '.lnk', '.url',  # 바로가기
        '.tmp', '.bak', '.swp',  # 임시 파일
        '.log',  # 로그 파일
        '.ini', '.cfg',  # 설정 파일
    }

    def index_directory(self, directory_path: str, progress_callback=None, settings: IndexingSettings = None) -> Dict:
        """Recursively index all files in a directory. 
        Documents get full content indexing, other files get metadata-only indexing."""
        logger.info(f"Starting index_directory: {directory_path}")
        print(f"Indexing directory: {directory_path}")
        self.errors = []  # Reset errors for each indexing session
        indexed_count = 0
        metadata_indexed_count = 0
        skipped_count = 0
        skipped_large = 0
        skipped_unchanged = 0  # Track files skipped due to no changes
        
        # Use default settings if not provided
        if settings is None:
            settings = IndexingSettings()
        
        logger.info(f"Settings: workers={settings.parallel_workers}, max_size={settings.max_file_size_mb}MB, skip_large={settings.skip_large_files}")
        
        # First pass: collect all files to index (all files, not just documents)
        files_to_index = []
        for root, dirs, files in os.walk(directory_path):
            dirs[:] = [d for d in dirs if not self._should_skip_dir(d, os.path.join(root, d))]
            for file in files:
                if self._should_skip_file(file):
                    continue
                ext = os.path.splitext(file)[1].lower()
                # Skip system/excluded extensions
                if ext in self.SKIP_EXTENSIONS:
                    continue
                file_path = os.path.join(root, file)
                files_to_index.append(file_path)
        
        total_files = len(files_to_index)
        processed_files = 0
        logger.info(f"Found {total_files} files to check")
        
        # Process files (sequential for now to avoid threading issues with embeddings)
        for file_path in files_to_index:
            file = os.path.basename(file_path)
            ext = os.path.splitext(file)[1].lower()
            is_content_indexable = ext in self.CONTENT_INDEXABLE_EXTS
            
            try:
                # Incremental indexing: Skip files that haven't changed
                if self.collection.is_file_indexed(file_path):
                    logger.debug(f"Skipped (unchanged): {file}")
                    skipped_unchanged += 1
                    skipped_count += 1
                    processed_files += 1
                    
                    # Still update progress for skipped files (every 5 files)
                    if progress_callback and total_files > 0:
                        if processed_files % 5 == 0:
                            percent = round((processed_files / total_files) * 100, 1)
                            progress_callback(f"⏭️ {file} (unchanged)", percent, processed_files, total_files)
                    continue
                
                # Check file size if skip_large_files is enabled (only for content-indexable files)
                if is_content_indexable and settings.skip_large_files:
                    file_size_mb = os.path.getsize(file_path) / (1024 * 1024)
                    if file_size_mb > settings.max_file_size_mb:
                        logger.debug(f"Skipped (large file {file_size_mb:.1f}MB): {file}")
                        skipped_large += 1
                        skipped_count += 1
                        processed_files += 1
                        continue
                
                # Index based on file type
                if is_content_indexable:
                    # Document files: full content indexing
                    logger.debug(f"Indexing content {processed_files+1}/{total_files}: {file}")
                    self.index_file(file_path, settings=settings)
                    indexed_count += 1
                else:
                    # Other files: metadata-only indexing
                    logger.debug(f"Indexing metadata {processed_files+1}/{total_files}: {file}")
                    self.index_metadata_only(file_path)
                    metadata_indexed_count += 1
                
                # Update file modification time after successful indexing
                self.collection.update_file_mtime(file_path)
                processed_files += 1
                
                # Log progress milestones
                if processed_files % 100 == 0:
                    logger.info(f"Progress: {processed_files}/{total_files} ({int(processed_files/total_files*100)}%) - Content: {indexed_count}, Metadata: {metadata_indexed_count}, Errors: {len(self.errors)}")
                
                # Call progress callback every 5 files for smoother updates
                if progress_callback and total_files > 0:
                    if processed_files % 5 == 0 or processed_files == total_files:
                        percent = round((processed_files / total_files) * 100, 1)
                        progress_callback(file, percent, processed_files, total_files)
                        
            except Exception as e:
                error_msg = f"Error indexing {file_path}: {e}"
                logger.error(error_msg, exc_info=True)
                self.errors.append({"file": file_path, "error": str(e)})
                processed_files += 1

        logger.info(f"Indexing complete. Content: {indexed_count}, Metadata: {metadata_indexed_count}, Skipped: {skipped_count} (Unchanged: {skipped_unchanged}, Large: {skipped_large}), Errors: {len(self.errors)}")
        return {
            "indexed": indexed_count + metadata_indexed_count,
            "indexed_content": indexed_count,
            "indexed_metadata": metadata_indexed_count,
            "skipped": skipped_count,
            "skipped_unchanged": skipped_unchanged,
            "skipped_large": skipped_large,
            "errors": self.errors
        }

        
    def index_outlook(self, max_emails: int = 1000) -> Dict:
        """Index emails using Outlook COM API (requires Outlook to be running)."""
        if not HAS_OUTLOOK:
            print("Outlook COM not available. Please install pywin32.")
            return {"indexed": 0, "skipped": 0, "errors": ["pywin32 not installed"]}
        
        print("Indexing Outlook emails via COM API...")
        indexed_count = 0
        skipped_count = 0
        self.errors = []
        
        try:
            pythoncom.CoInitialize()
            
            # Try to connect to running Outlook
            try:
                outlook = win32com.client.GetActiveObject("Outlook.Application")
                print("Connected to running Outlook instance.")
            except Exception:
                print("Outlook is not running. Please start Outlook first.")
                self.errors.append("Outlook is not running")
                return {"indexed": 0, "skipped": 0, "errors": self.errors}
            
            namespace = outlook.GetNamespace("MAPI")
            
            # Folders to index: Inbox(6), Sent Items(5)
            folders_to_index = [
                (6, "Inbox"),
                (5, "Sent Items")
            ]
            
            for folder_id, folder_name in folders_to_index:
                try:
                    folder = namespace.GetDefaultFolder(folder_id)
                    items = folder.Items
                    try:
                        items.Sort("[ReceivedTime]", True)  # Sort by date, newest first
                    except Exception as sort_error:
                        print(f"Warning: Could not sort items in {folder_name}: {sort_error}")

                    # Limit items per folder
                    items_per_folder = max_emails // len(folders_to_index)
                    count = 0
                    
                    print(f"Processing {folder_name}...")
                    
                    # Iterate with index to avoid iterator issues during modification
                    # But for simple reading, iteration is usually fine. 
                    # We'll stick to iteration but wrap item access carefully.
                    
                    # Note: accessing items directly in a loop can sometimes be flaky with COM
                    # if the collection changes. A safer way is to get a fixed list or handle errors gracefully.
                    
                    processed_items = 0
                    # We will try to iterate, but if it fails we just stop for this folder
                    try:
                        for item in items:
                            if count >= items_per_folder:
                                break
                            
                            processed_items += 1
                            if processed_items > items_per_folder * 2: # Fail-safe to prevent infinite loops if something is wrong
                                break

                            try:
                                # Only process mail items (Class 43)
                                # Check if item is valid object first
                                if item is None:
                                    continue
                                    
                                try:
                                    item_class = item.Class
                                except:
                                    # If we can't get Class, skip it
                                    continue

                                if item_class != 43:
                                    continue
                                
                                subject = getattr(item, 'Subject', "No Subject")
                                
                                sender = ""
                                try:
                                    sender = item.SenderName
                                    if not sender:
                                        sender = item.SenderEmailAddress
                                except:
                                    sender = "Unknown"
                                
                                # Get date
                                date = ""
                                try:
                                    date = item.ReceivedTime
                                except:
                                    try:
                                        date = item.SentOn
                                    except:
                                        pass
                                
                                # Get body (plain text preferred for indexing)
                                body = ""
                                try:
                                    body = item.Body
                                except:
                                    pass
                                
                                # Build full text for indexing
                                full_text = f"Subject: {subject}\nFrom: {sender}\nDate: {date}\n\n{body}"
                                
                                # Truncate body if too long (for embedding efficiency)
                                if len(full_text) > 10000:
                                    full_text = full_text[:10000]
                                
                                # Create display filename
                                safe_subject = "".join([c for c in subject if c.isalnum() or c in (' ', '-', '_', '[', ']')]).strip()[:60]
                                display_filename = f"[Email] {safe_subject}"
                                
                                # Create unique identifier using EntryID
                                try:
                                    entry_id = item.EntryID
                                except:
                                    import hashlib
                                    entry_id = hashlib.md5(f"{subject}{date}{sender}".encode()).hexdigest()
                                
                                # Chunking
                                chunks = self.text_splitter.split_text(full_text)
                                if not chunks:
                                    continue
                                
                                # Searchable header chunk (Korean + English keywords for better search)
                                header_chunk = f"파일명: {safe_subject} 파일: {display_filename} 발신인: {sender} 타입: 이메일 email outlook 메일"
                                chunks.insert(0, header_chunk)
                                
                                # Create unique IDs
                                import hashlib
                                id_hash = hashlib.md5(entry_id.encode() if isinstance(entry_id, str) else str(entry_id).encode()).hexdigest()
                                ids = [f"outlook_{id_hash}_{i}" for i in range(len(chunks))]
                                
                                # Metadata
                                metadatas = [{
                                    "source": "Outlook",
                                    "filename": display_filename,
                                    "created": str(date),
                                    "chunk_index": i,
                                    "type": "email",
                                    "path": f"outlook:{entry_id}",  # Special path for Outlook items
                                    "sender": sender,
                                    "folder": folder_name
                                } for i in range(len(chunks))]
                                
                                embeddings = self._get_embeddings(chunks)
                                
                                self.collection.upsert(
                                    documents=chunks,
                                    embeddings=embeddings,
                                    metadatas=metadatas,
                                    ids=ids
                                )
                                indexed_count += 1
                                count += 1
                                
                            except Exception as ex:
                                skipped_count += 1
                                # print(f"Error processing email object: {ex}")
                                continue
                                
                    except Exception as iter_error:
                        print(f"Error iterating items in {folder_name}: {iter_error}")
                        self.errors.append(f"Iteration error in {folder_name}: {iter_error}")

                    print(f"  {folder_name}: indexed {count} emails")
                    
                except Exception as folder_ex:
                    print(f"Error accessing {folder_name}: {folder_ex}")
                    self.errors.append(f"Folder error: {folder_name}")
            
        except Exception as e:
            print(f"Outlook COM Error: {e}")
            self.errors.append(str(e))
        finally:
            pythoncom.CoUninitialize()
        
        print(f"Outlook indexing complete. Indexed: {indexed_count}, Skipped: {skipped_count}")
        return {
            "indexed": indexed_count,
            "skipped": skipped_count,
            "errors": self.errors
        }

    # Phase 2: 3.1. FileIndexer에 추가해야 할 메서드
    def index_connector(self, connector) -> Dict:
        """
        Connector 기반 인덱싱 파이프라인.
        모든 문서는 동일한 인덱싱 로직을 통과하여 검색 품질에 차이가 없다.
        
        Args:
            connector: BaseConnector 구현체
            
        Returns:
            {"indexed": int, "skipped": int, "errors": list}
        """
        print(f"Indexing via connector: {connector.name}")
        self.errors = []
        indexed_count = 0
        skipped_count = 0
        
        # 인증 시도
        if not connector.authenticate():
            error_msg = f"Connector authentication failed: {connector.name}"
            print(error_msg)
            return {"indexed": 0, "skipped": 0, "errors": [error_msg]}
        
        try:
            for item in connector.list_items():
                try:
                    item_id = item.get("id", "")
                    source = item.get("source", "unknown")
                    text = item.get("text")
                    metadata = item.get("metadata", {})
                    filename = metadata.get("filename", "Unknown")
                    
                    # 텍스트가 직접 제공되지 않으면 파일에서 추출
                    if not text:
                        local_path = connector.download(item)
                        if local_path:
                            text = self.extract_text(local_path)
                        else:
                            skipped_count += 1
                            continue
                    
                    if not text or not text.strip():
                        skipped_count += 1
                        continue
                    
                    # 청킹
                    chunks = self.text_splitter.split_text(text)
                    if not chunks:
                        skipped_count += 1
                        continue
                    
                    # 파일명 청크 추가 (검색 품질 향상)
                    filename_no_ext = os.path.splitext(filename)[0] if filename else ""
                    filename_chunk = f"파일명: {filename_no_ext} 파일: {filename} 소스: {source}"
                    chunks.insert(0, filename_chunk)
                    
                    # ID 생성
                    import hashlib
                    id_hash = hashlib.md5(item_id.encode()).hexdigest()
                    ids = [f"{source}_{id_hash}_{i}" for i in range(len(chunks))]
                    
                    # 메타데이터 설정
                    chunk_metadatas = [{
                        "source": metadata.get("path") or item.get("path") or item_id,
                        "filename": filename,
                        "modified": metadata.get("modified"),
                        "author": metadata.get("author"),
                        "type": metadata.get("type", source),
                        "chunk_index": i,
                    } for i in range(len(chunks))]
                    
                    # 임베딩 생성
                    embeddings = self._get_embeddings(chunks)
                    
                    # 벡터 스토어에 저장
                    self.collection.upsert(
                        documents=chunks,
                        embeddings=embeddings,
                        metadatas=chunk_metadatas,
                        ids=ids
                    )
                    indexed_count += 1
                    
                except Exception as e:
                    error_msg = f"Error indexing item {item.get('id', 'unknown')}: {e}"
                    print(error_msg)
                    self.errors.append(error_msg)
                    skipped_count += 1
                    
        except Exception as e:
            error_msg = f"Connector iteration error: {e}"
            print(error_msg)
            self.errors.append(error_msg)
        finally:
            connector.close()
        
        print(f"Connector indexing complete. Indexed: {indexed_count}, Skipped: {skipped_count}")
        return {
            "indexed": indexed_count,
            "skipped": skipped_count,
            "errors": self.errors
        }

    def index_file(self, file_path: str, settings: IndexingSettings = None):
        """Process a single file: extract, chunk, and store vectors."""
        logger.debug(f"Processing: {file_path}")
        text = self.extract_text(file_path, settings=settings)
        
        filename = os.path.basename(file_path)
        filename_no_ext = os.path.splitext(filename)[0]
        
        # Prepend filename to content for better searchability
        # This allows searching by filename keywords
        filename_header = f"파일명: {filename_no_ext}\n파일: {filename}\n\n"
        text_with_filename = filename_header + text
        
        if not text_with_filename.strip():
            return

        chunks = self.text_splitter.split_text(text_with_filename)
        if not chunks:
            return
        
        # Also create a dedicated filename-only chunk for exact filename matching
        filename_chunk = f"파일명: {filename_no_ext} 파일: {filename} 경로: {file_path}"
        chunks.insert(0, filename_chunk)

        # Create unique IDs for chunks
        ids = [f"{file_path}_{i}" for i in range(len(chunks))]
        
        # Metadata for filtering and citation
        metadatas = [{
            "source": file_path,
            "filename": filename,
            "created": os.path.getctime(file_path),
            "chunk_index": i
        } for i in range(len(chunks))]

        # Generate embeddings manually
        embeddings = self._get_embeddings(chunks)

        # Upsert to vector store
        self.collection.upsert(
            documents=chunks,
            embeddings=embeddings,
            metadatas=metadatas,
            ids=ids
        )
        print(f"Indexed {len(chunks)} chunks for {file_path}")

    def index_metadata_only(self, file_path: str):
        """Index file metadata only (filename, extension, path) without content extraction.
        Used for non-document files like certificates, images, etc."""
        filename = os.path.basename(file_path)
        filename_no_ext = os.path.splitext(filename)[0]
        ext = os.path.splitext(filename)[1].lower()
        
        # Create searchable text from metadata
        # Include various variations for better search matching
        ext_name = ext.replace('.', '')  # e.g., "pfx", "p12"
        text = f"파일명: {filename_no_ext} 파일: {filename} 확장자: {ext} 형식: {ext_name} 경로: {file_path}"
        
        # Single chunk for metadata-only files
        ids = [f"{file_path}_meta"]
        
        try:
            created_time = os.path.getctime(file_path)
        except OSError:
            created_time = None
        
        metadatas = [{
            "source": file_path,
            "filename": filename,
            "created": created_time,
            "chunk_index": 0,
            "type": "metadata_only",  # Mark as metadata-only for UI differentiation
            "extension": ext
        }]
        
        embeddings = self._get_embeddings([text])
        
        self.collection.upsert(
            documents=[text],
            embeddings=embeddings,
            metadatas=metadatas,
            ids=ids
        )
        logger.debug(f"Indexed metadata for {filename}")


    def search(self, query: str, n_results: int = 5):
        """Search for relevant file chunks using hybrid search (semantic + keyword)."""
        if self.collection.count() == 0:
            return {
                'documents': [[]],
                'metadatas': [[]],
                'distances': [[]]
            }
        
        # Generate query embedding for semantic search
        query_embedding = self._get_embeddings([query])
        query_lower = query.lower()
        query_terms = query_lower.split()
        
        # Get more results for re-ranking
        initial_n = min(n_results * 4, self.collection.count())
        
        # Semantic search
        results = self.collection.query(
            query_embeddings=query_embedding,
            n_results=initial_n
        )
        
        if not results['documents'][0]:
            return results
        
        # Re-rank with keyword matching boost
        scored_results = []
        for i, doc in enumerate(results['documents'][0]):
            doc_lower = doc.lower()
            meta = results['metadatas'][0][i]
            filename_lower = meta.get('filename', '').lower()
            
            # Base score from semantic similarity (convert distance to similarity)
            semantic_score = 1 - results['distances'][0][i]
            
            # Keyword matching bonus (Reduced weights to prevent overriding semantic score)
            keyword_score = 0
            for term in query_terms:
                # Filename match (highest priority)
                if term in filename_lower:
                    keyword_score += 0.15
                # Content match
                if term in doc_lower:
                    keyword_score += 0.05
            
            # Combined score
            final_score = semantic_score + keyword_score
            
            scored_results.append({
                'doc': doc,
                'meta': meta,
                'distance': max(0, 1 - final_score),  # Convert back to distance
                'score': final_score
            })
        
        # Sort by score (higher is better)
        scored_results.sort(key=lambda x: x['score'], reverse=True)
        
        # Take top n results
        top_results = scored_results[:n_results]
        
        return {
            'documents': [[r['doc'] for r in top_results]],
            'metadatas': [[r['meta'] for r in top_results]],
            'distances': [[r['distance'] for r in top_results]]
        }


if __name__ == "__main__":
    # Test code
    indexer = FileIndexer()
    print("Indexer initialized. Ready to index files.")
    print(f"Current document count: {indexer.collection.count()}")
