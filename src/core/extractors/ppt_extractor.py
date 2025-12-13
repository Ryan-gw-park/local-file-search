"""
Local Finder X v2.0 - PowerPoint Extractor

Extracts content from .pptx files using python-pptx.
"""

from typing import List, Dict, Any

from .base import BaseExtractor, ExtractorResult, register_extractor

try:
    from pptx import Presentation
    from pptx.util import Inches
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    Presentation = None


@register_extractor
class PowerPointExtractor(BaseExtractor):
    """Extractor for Microsoft PowerPoint presentations (.pptx)."""
    
    SUPPORTED_EXTENSIONS = [".pptx"]
    
    def extract(self, file_path: str) -> ExtractorResult:
        """Extract content from a PowerPoint presentation."""
        if not PPTX_AVAILABLE:
            return self._create_error_result(
                "python-pptx is not installed. Install with: pip install python-pptx"
            )
        
        try:
            prs = Presentation(file_path)
            
            sections = []
            full_text_parts = []
            
            for slide_idx, slide in enumerate(prs.slides, 1):
                slide_content = self._extract_slide(slide, slide_idx)
                
                if slide_content["text"]:
                    sections.append({
                        "type": "slide",
                        "slide_number": slide_idx,
                        "title": slide_content.get("title", ""),
                        "content": slide_content["paragraphs"],
                    })
                    
                    # Format for full text
                    slide_header = f"## Slide {slide_idx}"
                    if slide_content.get("title"):
                        slide_header += f": {slide_content['title']}"
                    
                    full_text_parts.append(slide_header)
                    full_text_parts.append(slide_content["text"])
            
            full_text = "\n\n".join(full_text_parts)
            
            metadata = {
                "slide_count": len(prs.slides),
            }
            
            # Try to get presentation metadata
            try:
                if prs.core_properties.author:
                    metadata["author"] = prs.core_properties.author
                if prs.core_properties.title:
                    metadata["title"] = prs.core_properties.title
            except Exception:
                pass
            
            return self._create_success_result(
                text=full_text,
                sections=sections,
                metadata=metadata,
            )
            
        except Exception as e:
            return self._create_error_result(f"Error extracting PowerPoint: {str(e)}")
    
    def _extract_slide(self, slide, slide_idx: int) -> Dict[str, Any]:
        """Extract content from a single slide."""
        paragraphs = []
        title = None
        
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            # Check if this is the title
            if shape.is_placeholder:
                try:
                    if shape.placeholder_format.type == 1:  # Title placeholder
                        title = shape.text.strip()
                except Exception:
                    pass
            
            # Extract text from all text frames
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    paragraphs.append(text)
        
        # Also extract from tables if present
        for shape in slide.shapes:
            if shape.has_table:
                table_text = self._extract_table(shape.table)
                if table_text:
                    paragraphs.append(table_text)
        
        return {
            "title": title,
            "paragraphs": paragraphs,
            "text": "\n".join(paragraphs),
        }
    
    def _extract_table(self, table) -> str:
        """Extract text from a table."""
        rows = []
        for row in table.rows:
            cells = []
            for cell in row.cells:
                cells.append(cell.text.strip())
            rows.append(" | ".join(cells))
        return "\n".join(rows)


__all__ = ["PowerPointExtractor", "PPTX_AVAILABLE"]
